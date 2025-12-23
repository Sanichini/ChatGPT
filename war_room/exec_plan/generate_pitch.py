"""
generate_pitch.py
------------------

This script uses python‑pptx to generate a simple three‑slide
investor pitch deck for the new smartphone brand. The slides
highlight the market opportunity, product strategy, and financial
projections. Run this script to produce `investor_pitch.pptx` in
the current directory.

To run:

    python generate_pitch.py

Requires: python-pptx (listed in requirements.txt)
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def create_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = point
        p.font.size = Pt(16)
        p.level = 0
    return slide

def main():
    prs = Presentation()
    # Slide 1: Market opportunity
    slide1_points = [
        "Thailand smartphone market ~7.5 M units per quarter (Q3 2025)",
        "Mass‑market segment THB 5 k–15 k accounts for >40 % of shipments",
        "Top competitors: Samsung 24 %, OPPO 17 %, Xiaomi 15 %, Apple 13 %【693812817783168†L38-L53】",
        "5 G adoption and battery endurance drive consumer choice"
    ]
    create_slide(prs, "Market Opportunity & Landscape", slide1_points)
    # Slide 2: Product & Strategy
    slide2_points = [
        "Deliver flagship‑level features (120 Hz, Dimensity 8100, IP68) at THB 8 999",
        "Focus on Gen Z, young professionals and first‑time 5 G families",
        "Leverage real‑time analytics to time launch during competitor lulls",
        "Omni‑channel distribution: Shopee, Lazada, JD Central, Telco bundles"
    ]
    create_slide(prs, "Product & Go‑to‑Market Strategy", slide2_points)
    # Slide 3: Financials & KPIs
    slide3_points = [
        "Budget: THB 50 M (product dev 8 M, manufacturing 12 M, marketing 15 M, partners 5 M, support 2 M, contingency 8 M)",
        "Sales target: 100 k units in Q2 2025, ~1.3 % market share",
        "Optimal price point: THB 8 999 (forecast model) with predicted break‑even by Q4 2025",
        "KPIs: >80 % sell‑through in 60 days, 50 % brand awareness among target segments"
    ]
    create_slide(prs, "Financials & Key Metrics", slide3_points)
    prs.save('investor_pitch.pptx')

if __name__ == '__main__':
    main()